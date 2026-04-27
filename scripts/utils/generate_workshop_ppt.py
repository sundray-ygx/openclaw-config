#!/usr/bin/env python3
"""生成 AI-Native 落地推进研讨会 完整引导 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xD8, 0x3B, 0x01)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_PURPLE = RGBColor(0x6B, 0x5B, 0x95)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
TEXT_PRIMARY = RGBColor(0x33, 0x33, 0x33)
TEXT_SECONDARY = RGBColor(0x66, 0x66, 0x66)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
BG_ROW_ALT = RGBColor(0xF8, 0xF8, 0xF8)
BG_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_PRIMARY,
       align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def multi_line_tb(slide, left, top, width, height, lines, size=14, color=TEXT_PRIMARY, bold=False, line_spacing=0.35, align=PP_ALIGN.LEFT):
    """多行文本，每行一个paragraph"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(line_spacing * 20)
        p.alignment = align
    return txBox


def rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def rrect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def section_divider(title, subtitle="", time_label=""):
    """环节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_DARK)
    # 装饰线
    rect(slide, 1, 3.3, 1.8, 0.06, ACCENT_BLUE)
    tb(slide, 1, 2.4, 11, 0.9, title, 36, True, TEXT_WHITE)
    if subtitle:
        tb(slide, 1, 3.6, 11, 0.6, subtitle, 20, False, RGBColor(0xBB, 0xBB, 0xBB))
    if time_label:
        tb(slide, 1, 6, 5, 0.4, time_label, 14, False, RGBColor(0x88, 0x88, 0x88))
    return slide


def content_slide(title):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_LIGHT)
    rect(slide, 0, 0, 13.333, 0.06, ACCENT_BLUE)
    tb(slide, 0.8, 0.3, 11.5, 0.7, title, 28, True, TEXT_PRIMARY)
    rect(slide, 0.8, 1.05, 11.5, 0.015, BORDER_GRAY)
    return slide


def table_row(slide, x, y, cols, bg=TEXT_WHITE, border=BORDER_GRAY):
    """画一行表格"""
    for (cx, cw, text, size, bold, color) in cols:
        rect(slide, cx, y, cw, 0.5, bg, border, 0.5)
        tb(slide, cx + 0.08, y + 0.05, cw - 0.16, 0.4, text, size, bold, color)


# ══════════════════════════════════════════════════
# PART A: 开场共识引导（原有7页）
# ══════════════════════════════════════════════════

# --- Slide 1: 封面 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
rect(slide, 1, 3.2, 2, 0.06, ACCENT_BLUE)
tb(slide, 1, 2, 11, 1.2, "AI-Native 落地推进研讨会", 36, True, TEXT_WHITE)
tb(slide, 1, 3.5, 11, 0.8, "对齐障碍 · 拆解目标 · 建立检视", 20, False, RGBColor(0xAA, 0xAA, 0xAA))
tb(slide, 1, 6.2, 11, 0.5, "2026年4月27日  19:00 - 21:30 ｜ 参会：无线 · 交换机 · NMC · 会议主机 · IPSIP · 硬件", 14, False, RGBColor(0x88, 0x88, 0x88))

# --- Slide 2: 议程总览 ---
slide = content_slide("今日议程")

agenda = [
    ("19:00", "19:15", "开场：方法论共识", "统一方法论认知，确认研讨规则", ACCENT_BLUE, "15min"),
    ("19:15", "19:55", "环节一：障碍对齐", "各产线TOP3障碍呈现 + 集体诊断分类", ACCENT_RED, "40min"),
    ("19:55", "20:05", "☕ 休息", "", TEXT_SECONDARY, "10min"),
    ("20:05", "20:55", "环节二：目标拆解", "近期里程碑SMART验证 + 具体行动项拆解", ACCENT_ORANGE, "50min"),
    ("20:55", "21:30", "环节三：衡量与共识", "衡量标准 + 检视节奏 + 行动项确认", ACCENT_GREEN, "35min"),
]

# 表头
for name, w, x in [("时间", 1.5, 0.8), ("环节", 3, 2.5), ("内容", 4.5, 5.7), ("时长", 1, 10.4)]:
    rect(slide, x, 1.3, w, 0.5, RGBColor(0x33, 0x33, 0x33))
    tb(slide, x, 1.33, w, 0.44, name, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)

for i, (start, end, title, desc, color, dur) in enumerate(agenda):
    y = 1.85 + i * 0.9
    bg = TEXT_WHITE if i % 2 == 0 else BG_ROW_ALT
    is_break = title.startswith("☕")
    c = color if not is_break else TEXT_SECONDARY

    # 时间
    rect(slide, 0.8, y, 1.5, 0.75, bg, BORDER_GRAY, 0.5)
    tb(slide, 0.85, y+0.1, 1.4, 0.55, f"{start}\n{end}", 13, True, c, PP_ALIGN.CENTER)
    # 环节
    rect(slide, 2.5, y, 3, 0.75, bg, BORDER_GRAY, 0.5)
    if not is_break:
        rrect(slide, 2.7, y+0.15, 2.6, 0.45, color)
        tb(slide, 2.7, y+0.18, 2.6, 0.4, title, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)
    else:
        tb(slide, 2.7, y+0.15, 2.6, 0.45, title, 14, False, TEXT_SECONDARY, PP_ALIGN.CENTER)
    # 内容
    rect(slide, 5.7, y, 4.5, 0.75, bg, BORDER_GRAY, 0.5)
    tb(slide, 5.9, y+0.1, 4.2, 0.55, desc, 13, False, TEXT_PRIMARY)
    # 时长
    rect(slide, 10.4, y, 1, 0.75, bg, BORDER_GRAY, 0.5)
    tb(slide, 10.4, y+0.15, 1, 0.4, dur, 15, True, c, PP_ALIGN.CENTER)

# 底部参会信息
rrect(slide, 0.8, 6.5, 11.5, 0.6, BG_LIGHT_BLUE, ACCENT_BLUE, 1)
tb(slide, 1, 6.55, 11, 0.5, "参会产线：无线 · 交换机 · NMC · 会议主机 · IPSIP · 硬件 ｜ 核心原则：解题会，不是汇报会", 14, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# --- Slide 3: 我们在哪 ---
slide = content_slide("一、我们在哪：当前状态")

tb(slide, 0.8, 1.3, 5, 0.4, "✅ 已完成的工作", 18, True, ACCENT_GREEN)
done = ["各产线有年度OKR和里程碑计划", "综合管理部对齐报告（V1.2）",
        "共性需求已识别（统一平台、知识库、度量、流程规范）", "里程碑汇总表（V2.0）"]
for i, item in enumerate(done):
    tb(slide, 1.2, 1.8 + i * 0.38, 4.5, 0.35, "• " + item, 14, False, TEXT_PRIMARY)

tb(slide, 6.5, 1.3, 5.5, 0.4, "⚠️ 当前存在的问题", 18, True, ACCENT_ORANGE)
problems = [("障碍卡住没解", "缺乏系统性解决路径，进度停滞"),
            ("目标太大太远", "缺乏具体可执行的任务拆解"),
            ("没有过程衡量", "年底才发现偏差，来不及调整")]
for i, (t, d) in enumerate(problems):
    y = 1.85 + i * 0.7
    tb(slide, 6.8, y, 5.5, 0.3, "🔴 " + t, 14, True, ACCENT_RED)
    tb(slide, 7.3, y + 0.3, 5.2, 0.3, d, 12, False, TEXT_SECONDARY)

rrect(slide, 1.5, 5.8, 10.3, 1, ACCENT_BLUE)
tb(slide, 1.5, 5.85, 10.3, 0.5, "今天的定位：不是汇报会，是解题会", 22, True, TEXT_WHITE, PP_ALIGN.CENTER)
tb(slide, 1.5, 6.35, 10.3, 0.4, "把「我们要去哪里」变成「下一步具体干什么、怎么知道干对了」", 15, False, TEXT_WHITE, PP_ALIGN.CENTER)

# --- Slide 4: 方法论框架 ---
slide = content_slide("二、方法论框架：从战略到执行")

# 战略层卡片
rrect(slide, 1, 1.5, 4.5, 3.2, TEXT_WHITE, BORDER_GRAY, 1)
tb(slide, 1.2, 1.6, 4, 0.4, "战略层（已完成）", 16, True, ACCENT_GREEN)
multi_line_tb(slide, 1.2, 2.1, 4, 2.2, [
    "RPP — 优先事项 / 流程 / 资源",
    "PPT — 人 / 流程 / 技术",
    "",
    "解决「为什么做、怎么组织」",
    "已有基础，不需要重复讨论"
], 13, TEXT_PRIMARY)

# 箭头
tb(slide, 5.7, 2.7, 0.8, 0.5, "拆解", 14, True, ACCENT_BLUE, PP_ALIGN.CENTER)
tb(slide, 5.7, 3.05, 0.8, 0.5, "──→", 24, True, ACCENT_BLUE, PP_ALIGN.CENTER)

# 执行层卡片
rrect(slide, 6.8, 1.5, 5.5, 3.2, TEXT_WHITE, BORDER_GRAY, 1)
tb(slide, 7, 1.6, 5, 0.4, "执行层（今天引入）", 16, True, ACCENT_BLUE)
multi_line_tb(slide, 7, 2.1, 5, 2.2, [
    "PGSAR闭环",
    "  P(Plan) → G(Goal) → S(Standard)",
    "  → A(Action) → R(Review) → 循环",
    "",
    "解决「具体做什么、怎么衡量、怎么纠偏」"
], 13, TEXT_PRIMARY)

# SMART
rrect(slide, 3, 5.2, 7.3, 0.9, BG_LIGHT_BLUE, ACCENT_BLUE, 1)
tb(slide, 3.2, 5.3, 7, 0.4, "SMART原则 — 目标质量检查工具", 16, True, ACCENT_BLUE, PP_ALIGN.CENTER)
tb(slide, 3.2, 5.7, 7, 0.3, "Specific · Measurable · Achievable · Relevant · Time-bound", 14, False, TEXT_PRIMARY, PP_ALIGN.CENTER)

# --- Slide 5: PGSAR闭环 ---
slide = content_slide("三、PGSAR 执行闭环")

steps = [("P", "Plan", "检视整体计划\n确认方向", ACCENT_BLUE),
         ("G", "Goal", "SMART拆解\n大目标→小目标", ACCENT_PURPLE),
         ("S", "Standard", "定义「做好了」\n的标准，可衡量", ACCENT_TEAL),
         ("A", "Action", "明确行动\n责任人、时间", ACCENT_ORANGE),
         ("R", "Review", "定期检视\n红黄绿灯纠偏", ACCENT_RED)]

bw, bh, gap = 2.1, 1.3, 0.35
x_start = 1
for i, (letter, name, desc, color) in enumerate(steps):
    x = x_start + i * (bw + gap)
    rrect(slide, x, 1.5, bw, bh, color)
    tb(slide, x, 1.6, bw, 0.4, f"{letter} - {name}", 17, True, TEXT_WHITE, PP_ALIGN.CENTER)
    tb(slide, x, 2.05, bw, 0.7, desc, 12, False, TEXT_WHITE, PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        tb(slide, x + bw - 0.05, 1.9, 0.45, 0.4, "→", 22, True, color, PP_ALIGN.CENTER)

tb(slide, 8.5, 3, 3, 0.4, "↑ 回到Plan，持续循环", 14, False, TEXT_SECONDARY, PP_ALIGN.CENTER)

rrect(slide, 2, 3.6, 9.3, 0.7, BG_LIGHT_BLUE)
tb(slide, 2, 3.65, 9.3, 0.6, "计划 → 拆解 → 定标准 → 行动 → 检视 → 调整，循环往复", 18, True, ACCENT_BLUE, PP_ALIGN.CENTER)

tb(slide, 0.8, 4.8, 11, 0.4, "本次研讨会在PGSAR中的位置：", 16, True, TEXT_PRIMARY)

focus = [("障碍对齐", "检视Plan，识别阻塞项", ACCENT_RED),
         ("目标拆解", "Goal拆解 + SMART验证", ACCENT_BLUE),
         ("建立检视", "Standard + Review机制", ACCENT_GREEN)]
for i, (t, d, c) in enumerate(focus):
    x = 1 + i * 4
    rrect(slide, x, 5.3, 3.5, 0.9, TEXT_WHITE, c, 2)
    tb(slide, x + 0.15, 5.35, 3.2, 0.35, t, 16, True, c, PP_ALIGN.CENTER)
    tb(slide, x + 0.15, 5.7, 3.2, 0.35, d, 13, False, TEXT_SECONDARY, PP_ALIGN.CENTER)

# --- Slide 6: SMART原则 ---
slide = content_slide("四、SMART原则 — 目标质量检查")

headers = ["维度", "问什么", "✅ 合格", "❌ 不合格"]
hx = [0.8, 2.5, 5.5, 9]
hw = [1.5, 2.8, 3.3, 3.5]
for h, x, w in zip(headers, hx, hw):
    rect(slide, x, 1.4, w, 0.45, RGBColor(0x33, 0x33, 0x33))
    tb(slide, x, 1.42, w, 0.4, h, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)

smart = [("S", "Specific", "具体做什么？交付什么？", "有明确交付物和范围", "\"推进AI应用\"", ACCENT_BLUE),
         ("M", "Measurable", "怎么衡量完成？", "有量化指标和基线", "\"提升效率\"", ACCENT_PURPLE),
         ("A", "Achievable", "现在能做吗？", "信心指数≥7", "目标远大无路径", ACCENT_ORANGE),
         ("R", "Relevant", "和OKR什么关系？", "能追溯到具体KR", "孤立任务", ACCENT_TEAL),
         ("T", "Time-bound", "什么时候完成？", "有明确截止日和节点", "\"尽快\"\"年底前\"", ACCENT_RED)]

for ri, (l, n, q, good, bad, c) in enumerate(smart):
    y = 1.9 + ri * 0.6
    bg = TEXT_WHITE if ri % 2 == 0 else BG_ROW_ALT
    for x, w in zip(hx, hw):
        rect(slide, x, y, w, 0.55, bg, BORDER_GRAY, 0.5)
    tb(slide, hx[0]+0.1, y+0.07, hw[0]-0.2, 0.4, f"{l} {n}", 12, True, c)
    tb(slide, hx[1]+0.1, y+0.07, hw[1]-0.2, 0.4, q, 13, False, TEXT_PRIMARY)
    tb(slide, hx[2]+0.1, y+0.07, hw[2]-0.2, 0.4, good, 13, False, ACCENT_GREEN)
    tb(slide, hx[3]+0.1, y+0.07, hw[3]-0.2, 0.4, bad, 13, False, ACCENT_RED)

tb(slide, 0.8, 5.2, 11, 0.4, "核心工具：信心指数  __ / 10", 20, True, TEXT_PRIMARY)

levels = [("≥8 分", "目标合理，直接拆解行动", ACCENT_GREEN),
          ("5-7 分", "有风险，讨论卡在哪，可能要拆分", ACCENT_ORANGE),
          ("<5 分", "重新定义或大幅拆分", ACCENT_RED)]
for i, (s, d, c) in enumerate(levels):
    x = 0.8 + i * 4.2
    rrect(slide, x, 5.75, 3.8, 0.75, TEXT_WHITE, c, 2)
    tb(slide, x+0.15, 5.8, 1.2, 0.3, s, 16, True, c)
    tb(slide, x+0.15, 6.1, 3.5, 0.3, d, 12, False, TEXT_SECONDARY)

# --- Slide 7: 三个产出 ---
slide = content_slide("五、今天的三个产出")

outputs = [("1", "障碍清单", "各产线TOP3障碍\n分类 · 优先级 · 解法方向", "知道卡在哪、谁来解", ACCENT_RED),
           ("2", "目标拆解", "近期里程碑\nSMART验证 · 行动项", "知道下一步干什么", ACCENT_BLUE),
           ("3", "检视机制", "过程指标\n红黄绿灯 · 检视节奏", "知道做得怎样、及时纠偏", ACCENT_GREEN)]

for i, (num, title, desc, value, color) in enumerate(outputs):
    x = 1 + i * 4
    rrect(slide, x, 1.6, 3.5, 4, TEXT_WHITE, color, 2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+1.3), Inches(1.9), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tb(slide, x+1.3, 2.0, 0.8, 0.8, num, 28, True, TEXT_WHITE, PP_ALIGN.CENTER)
    tb(slide, x+0.3, 2.9, 2.9, 0.5, title, 24, True, color, PP_ALIGN.CENTER)
    tb(slide, x+0.3, 3.5, 2.9, 1.2, desc, 14, False, TEXT_PRIMARY, PP_ALIGN.CENTER)
    tb(slide, x+0.3, 4.9, 2.9, 0.4, value, 13, True, TEXT_SECONDARY, PP_ALIGN.CENTER)

# --- Slide 8: 共识确认 ---
slide = content_slide("六、共识确认")

tb(slide, 0.8, 1.2, 11, 0.4, "请逐条确认，有异议现场提出：", 16, False, TEXT_SECONDARY)

consensus = [("共识一", "今天是解题会，不是汇报会 — 重点在「怎么解决」", ACCENT_BLUE),
             ("共识二", "用SMART验证目标 — 信心<7的目标必须讨论拆分", ACCENT_PURPLE),
             ("共识三", "每个行动有唯一责任人 — 不能写「团队负责」", ACCENT_ORANGE),
             ("共识四", "建立双周红黄绿灯检视 — 不等到年底才发现偏差", ACCENT_TEAL),
             ("共识五", "诚实暴露问题 — 卡住了就说卡住了", ACCENT_RED)]

for i, (t, d, c) in enumerate(consensus):
    y = 1.9 + i * 0.85
    rrect(slide, 1, y, 1.5, 0.55, c)
    tb(slide, 1, y+0.07, 1.5, 0.4, t, 16, True, TEXT_WHITE, PP_ALIGN.CENTER)
    rrect(slide, 2.7, y, 9, 0.55, TEXT_WHITE, BORDER_GRAY, 1)
    tb(slide, 2.9, y+0.07, 8.6, 0.4, d, 16, False, TEXT_PRIMARY)

tb(slide, 2, 6.5, 9, 0.5, "✅ 确认完毕 → 进入环节一：障碍对齐", 18, True, ACCENT_BLUE, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# PART B: 环节一 — 障碍对齐
# ══════════════════════════════════════════════════

section_divider("环节一：障碍对齐", "各产线TOP3障碍 + 集体诊断", "19:15 - 19:55（40分钟）")

# --- 障碍对齐规则 ---
slide = content_slide("环节一：障碍对齐 — 规则")

tb(slide, 0.8, 1.3, 11, 0.4, "各产线轮流呈现，每产线 5-6 分钟", 20, True, ACCENT_BLUE)

rules = ["每个产线讲 TOP3 障碍：是什么、影响多大、已尝试什么、需要什么支持",
         "其他产线可补充类似经验，但不展开讨论（避免超时）",
         "主持人现场分类，标记优先级",
         "目标是收集和分类，不是现场解决所有问题"]
for i, r in enumerate(rules):
    y = 2.1 + i * 0.5
    rrect(slide, 1, y, 11, 0.42, TEXT_WHITE, BORDER_GRAY, 1)
    tb(slide, 1.2, y+0.03, 10.5, 0.36, f"{'❶❷❸❹'[i]}  {r}", 15, False, TEXT_PRIMARY)

# 每个产线时间分配
tb(slide, 0.8, 4.3, 11, 0.4, "时间分配：6个产线 × 5分钟 + 聚类 10分钟 = 40分钟", 16, True, TEXT_PRIMARY)

teams = ["无线", "交换机", "NMC", "会议主机", "IPSIP", "硬件"]
for i, team in enumerate(teams):
    x = 0.8 + i * 1.9
    c = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_TEAL, ACCENT_ORANGE, RGBColor(0x7B, 0x83, 0xEB), RGBColor(0xE9, 0x7C, 0x23)][i]
    rrect(slide, x, 4.85, 1.7, 0.55, TEXT_WHITE, c, 1.5)
    tb(slide, x, 4.9, 1.7, 0.4, team, 13, True, c, PP_ALIGN.CENTER)

tb(slide, 0.8, 5.7, 11, 0.35, "之后 10分钟 → 跨产线共性障碍聚类 + 优先级排序", 15, True, ACCENT_ORANGE)

# --- 障碍分类框架 ---
slide = content_slide("环节一：障碍分类框架")

tb(slide, 0.8, 1.2, 11, 0.4, "主持人现场分类到以下5个类别：", 16, False, TEXT_SECONDARY)

cats = [
    ("A", "工具/平台缺失", "基础设施不到位，想用但没工具", "加速建设、找临时替代", ACCENT_BLUE),
    ("B", "能力/技能不足", "工具在但不会用、用不好", "培训、布道师辅导、结对", ACCENT_PURPLE),
    ("C", "流程/规范阻力", "现有流程与AI-Native冲突", "流程适配、试点验证", ACCENT_ORANGE),
    ("D", "资源/优先级冲突", "人不够、和其他工作排不开", "资源协调、管理层决策", ACCENT_RED),
    ("E", "目标不清晰", "不知道具体该做什么、怎么做", "目标拆解（环节二解决）", ACCENT_TEAL),
]

for i, (letter, name, desc, solution, color) in enumerate(cats):
    y = 1.8 + i * 1.0
    # 标签
    rrect(slide, 0.8, y, 0.8, 0.7, color)
    tb(slide, 0.8, y+0.1, 0.8, 0.5, letter, 24, True, TEXT_WHITE, PP_ALIGN.CENTER)
    # 名称
    tb(slide, 1.8, y+0.02, 2.5, 0.35, name, 17, True, color)
    # 描述
    tb(slide, 4.5, y+0.02, 4, 0.35, desc, 14, False, TEXT_PRIMARY)
    # 解法
    rrect(slide, 8.8, y+0.05, 3.8, 0.6, BG_LIGHT_BLUE, color, 1)
    tb(slide, 9, y+0.1, 3.5, 0.5, solution, 13, False, color)

tb(slide, 0.8, 6.8, 11, 0.3, "标注「2+产线共同面临」的障碍 → 提升为平台级问题，综合管理部优先处理", 14, True, ACCENT_RED)


# ══════════════════════════════════════════════════
# PART C: 环节二 — 目标拆解工作坊
# ══════════════════════════════════════════════════

section_divider("环节二：目标拆解工作坊", "近期里程碑 SMART验证 + 任务拆解", "20:05 - 20:55（50分钟）")

# --- 时间分配 ---
slide = content_slide("环节二：时间分配")

tb(slide, 0.8, 1.3, 11, 0.4, "6个产线 × 7分钟/产线 + 缓冲 8分钟 = 50分钟", 20, True, ACCENT_BLUE)

team_times = [
    ("无线", "20:05", "20:12", "1-2个近期里程碑\nSMART验证 + 行动项", ACCENT_BLUE),
    ("交换机", "20:12", "20:19", "1-2个近期里程碑\nSMART验证 + 行动项", ACCENT_PURPLE),
    ("NMC", "20:19", "20:26", "1-2个近期里程碑\nSMART验证 + 行动项", ACCENT_TEAL),
    ("会议主机", "20:26", "20:33", "1-2个近期里程碑\nSMART验证 + 行动项", ACCENT_ORANGE),
    ("IPSIP", "20:33", "20:40", "1-2个近期目标\nSMART验证 + 行动项", RGBColor(0x7B, 0x83, 0xEB)),
    ("硬件", "20:40", "20:47", "1-2个近期目标\nSMART验证 + 行动项", RGBColor(0xE9, 0x7C, 0x23)),
]

# 表头
for name, w, x in [("产线", 1.8, 0.8), ("时间", 1.5, 2.8), ("内容", 3.5, 4.5), ("状态", 1.2, 8.2)]:
    rect(slide, x, 2, w, 0.45, RGBColor(0x33, 0x33, 0x33))
    tb(slide, x, 2.03, w, 0.4, name, 13, True, TEXT_WHITE, PP_ALIGN.CENTER)

for i, (team, start, end, content, color) in enumerate(team_times):
    y = 2.5 + i * 0.8
    bg = TEXT_WHITE if i % 2 == 0 else BG_ROW_ALT
    rect(slide, 0.8, y, 1.8, 0.7, bg, BORDER_GRAY, 0.5)
    rrect(slide, 0.95, y+0.12, 1.5, 0.45, color)
    tb(slide, 0.95, y+0.15, 1.5, 0.4, team, 15, True, TEXT_WHITE, PP_ALIGN.CENTER)
    rect(slide, 2.8, y, 1.5, 0.7, bg, BORDER_GRAY, 0.5)
    tb(slide, 2.85, y+0.1, 1.4, 0.5, f"{start}\n{end}", 12, True, TEXT_PRIMARY, PP_ALIGN.CENTER)
    rect(slide, 4.5, y, 3.5, 0.7, bg, BORDER_GRAY, 0.5)
    tb(slide, 4.7, y+0.05, 3.2, 0.6, content, 12, False, TEXT_PRIMARY)
    rect(slide, 8.2, y, 1.2, 0.7, bg, BORDER_GRAY, 0.5)
    tb(slide, 8.2, y+0.15, 1.2, 0.4, "⏳", 16, False, TEXT_SECONDARY, PP_ALIGN.CENTER)

# 缓冲说明
rrect(slide, 0.8, 6.6, 11.5, 0.5, BG_LIGHT_BLUE, ACCENT_ORANGE, 1)
tb(slide, 1, 6.65, 11, 0.4, "20:47 - 20:55  缓冲时间：补充讨论未完成的拆解 / 回答共性问题", 14, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

# --- SMART拆解模板 ---
slide = content_slide("环节二：SMART拆解模板")

# 左侧模板
tb(slide, 0.8, 1.3, 5.5, 0.4, "里程碑拆解模板（每产线1-2个近期里程碑）", 17, True, ACCENT_BLUE)

template_items = [
    ("里程碑名称：", "________________", 0.8, 1.85),
    ("信心指数：", "___ / 10", 0.8, 2.25),
]

for label, value, x, y in template_items:
    tb(slide, x, y, 2, 0.35, label, 14, True, TEXT_PRIMARY)
    tb(slide, x+2, y, 3.5, 0.35, value, 14, False, TEXT_SECONDARY)

# SMART表格
smart_labels = ["S 具体做什么？交付什么？", "M 怎么衡量完成？", "A 现在能做吗？卡在哪？",
                "R 和年度OKR什么关系？", "T 什么时候完成？"]
for i, label in enumerate(smart_labels):
    y = 2.8 + i * 0.5
    colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_TEAL, ACCENT_RED]
    rect(slide, 0.8, y, 1.8, 0.42, colors[i])
    tb(slide, 0.85, y+0.03, 1.7, 0.36, label.split(" ")[0], 13, True, TEXT_WHITE, PP_ALIGN.CENTER)
    rect(slide, 2.6, y, 4.2, 0.42, TEXT_WHITE, BORDER_GRAY, 0.5)
    tb(slide, 2.7, y+0.03, 4, 0.36, label.split(" ", 1)[1], 13, False, TEXT_PRIMARY)

# 右侧引导问题
tb(slide, 7.5, 1.3, 5, 0.4, "主持人引导问题", 17, True, ACCENT_ORANGE)

questions = [
    ("目标模糊时", "做完怎么知道做成了？\n看到什么就算完成？"),
    ("信心指数低时", "为什么不是8或9？\n最不确定的是什么？"),
    ("行动不具体时", "如果只做一件事推进这个目标，\n你做什么？"),
    ("责任人不清晰时", "谁来做？什么时候能看到\n初步结果？"),
    ("目标太大时", "切成两半，前半段是什么？"),
    ("和OKR关联不清时", "这个目标达成后，\n哪个KR会怎么变化？"),
]

for i, (situation, question) in enumerate(questions):
    y = 1.85 + i * 0.85
    rrect(slide, 7.5, y, 5, 0.75, TEXT_WHITE, ACCENT_ORANGE, 1)
    tb(slide, 7.7, y+0.03, 4.6, 0.3, situation, 13, True, ACCENT_ORANGE)
    tb(slide, 7.7, y+0.33, 4.6, 0.4, question, 12, False, TEXT_PRIMARY)

# 底部拆解深度要求
tb(slide, 0.8, 5.6, 11, 0.35, "拆解深度要求：", 15, True, TEXT_PRIMARY)
rrect(slide, 0.8, 6.0, 5.5, 0.5, BG_LIGHT_BLUE, ACCENT_GREEN, 1)
tb(slide, 1, 6.05, 5, 0.4, "✅ 有具体行动 + 唯一责任人 + 时间 + 交付物", 13, False, ACCENT_GREEN)
rrect(slide, 7, 6.0, 5.5, 0.5, RGBColor(0xFD, 0xE8, 0xE8), ACCENT_RED, 1)
tb(slide, 7.2, 6.05, 5, 0.4, "❌ 「推进XX」「加强XX」「持续优化XX」", 13, False, ACCENT_RED)

# --- 行动项表格 ---
slide = content_slide("环节二：行动项格式")

tb(slide, 0.8, 1.3, 11, 0.4, "每个里程碑拆出的行动项必须满足以下格式：", 16, False, TEXT_SECONDARY)

# 表头
cols_h = [("#", 0.5, 0.8), ("具体行动", 3.5, 1.5), ("责任人", 1.2, 5.3),
          ("开始时间", 1.2, 6.5), ("完成时间", 1.2, 7.7), ("交付物", 2, 8.9)]
for name, w, x in cols_h:
    rect(slide, x, 1.9, w, 0.5, RGBColor(0x33, 0x33, 0x33))
    tb(slide, x, 1.93, w, 0.44, name, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)

# 示例行
example_rows = [
    ("1", "日志分析skill标准化开发", "张三", "4/28", "5/10", "skill代码+文档"),
    ("2", "团队skill使用培训", "李四", "5/5", "5/12", "培训完成率≥90%"),
    ("3", "知识库框架搭建", "王五", "5/8", "5/20", "知识库框架+数据结构"),
]
for ri, row in enumerate(example_rows):
    y = 2.4 + ri * 0.5
    bg = TEXT_WHITE if ri % 2 == 0 else BG_ROW_ALT
    for ci, (val, w, x) in enumerate(zip(row, [c[1] for c in cols_h], [c[2] for c in cols_h])):
        rect(slide, x, y, w, 0.45, bg, BORDER_GRAY, 0.5)
        tb(slide, x+0.05, y+0.03, w-0.1, 0.4, val, 12, False, TEXT_PRIMARY if ci > 0 else TEXT_SECONDARY, PP_ALIGN.CENTER if ci != 1 else PP_ALIGN.LEFT)

# 关键规则
tb(slide, 0.8, 4.2, 11, 0.4, "关键规则：", 16, True, TEXT_PRIMARY)

rules = [
    ("唯一责任人", "必须落实到具体个人，不能写「团队」「待定」", ACCENT_RED),
    ("具体交付物", "每个行动必须有明确的产出物，不能只写「完成」", ACCENT_BLUE),
    ("可验证时间", "有明确的完成日期，不能写「尽快」「Q2」", ACCENT_ORANGE),
]
for i, (title, desc, c) in enumerate(rules):
    y = 4.7 + i * 0.65
    rrect(slide, 1, y, 1.8, 0.5, c)
    tb(slide, 1, y+0.05, 1.8, 0.4, title, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)
    rrect(slide, 3, y, 9, 0.5, TEXT_WHITE, BORDER_GRAY, 1)
    tb(slide, 3.2, y+0.05, 8.6, 0.4, desc, 14, False, TEXT_PRIMARY)


# ══════════════════════════════════════════════════
# PART D: 环节三 — 衡量标准与共识确认
# ══════════════════════════════════════════════════

section_divider("环节三：衡量标准与共识确认", "衡量标准 + 检视节奏 + 行动项确认", "20:55 - 21:30（35分钟）")

# --- 衡量标准 ---
slide = content_slide("环节三：衡量标准设计")

tb(slide, 0.8, 1.2, 5, 0.4, "前置指标（Leading）— 过程中持续关注", 17, True, ACCENT_BLUE)
tb(slide, 0.8, 1.6, 5, 0.3, "提前预警，发现问题于未然", 13, False, TEXT_SECONDARY)

# 前置指标示例
lead_h = [("指标", 3, 0.8), ("含义", 2.5, 3.8), ("目标值", 1.5, 6.3), ("采集频率", 1.2, 7.8)]
for name, w, x in lead_h:
    rect(slide, x, 2, w, 0.42, ACCENT_BLUE)
    tb(slide, x, 2.02, w, 0.38, name, 12, True, TEXT_WHITE, PP_ALIGN.CENTER)

lead_data = [
    ("AI工具日活率", "团队成员日常使用AI工具的比例", "≥80%", "每周"),
    ("SPEC文档覆盖率", "需求/设计阶段SPEC文档使用比例", "≥60%", "双周"),
    ("CI通过率", "代码提交后CI自动通过的比例", "≥90%", "每日"),
    ("技能使用频次", "各Skill被调用的频次", "持续增长", "每周"),
]
for ri, (metric, meaning, target, freq) in enumerate(lead_data):
    y = 2.42 + ri * 0.45
    bg = TEXT_WHITE if ri % 2 == 0 else BG_ROW_ALT
    for val, (_, w, x) in zip([metric, meaning, target, freq], lead_h):
        rect(slide, x, y, w, 0.42, bg, BORDER_GRAY, 0.5)
        tb(slide, x+0.05, y+0.03, w-0.1, 0.36, val, 11, False, TEXT_PRIMARY, PP_ALIGN.CENTER)

# 滞后指标
tb(slide, 0.8, 4.4, 5, 0.4, "滞后指标（Lagging）— 阶段末衡量结果", 17, True, ACCENT_ORANGE)
tb(slide, 0.8, 4.8, 5, 0.3, "验证结果，评估目标达成", 13, False, TEXT_SECONDARY)

lag_h = [("指标", 3, 0.8), ("含义", 2.5, 3.8), ("目标值", 1.5, 6.3), ("衡量时间", 1.2, 7.8)]
for name, w, x in lag_h:
    rect(slide, x, 5.15, w, 0.42, ACCENT_ORANGE)
    tb(slide, x, 5.17, w, 0.38, name, 12, True, TEXT_WHITE, PP_ALIGN.CENTER)

lag_data = [
    ("需求交付周期", "需求确认到交付的端到端时间", "缩短40-50%", "季末"),
    ("技术债务消除率", "已清理TD占总TD的比例", "≥80%", "季末"),
    ("项目周期缩短率", "对比基线的项目周期变化", "缩短25-40%", "季末"),
]
for ri, (metric, meaning, target, freq) in enumerate(lag_data):
    y = 5.57 + ri * 0.42
    bg = TEXT_WHITE if ri % 2 == 0 else BG_ROW_ALT
    for val, (_, w, x) in zip([metric, meaning, target, freq], lag_h):
        rect(slide, x, y, w, 0.4, bg, BORDER_GRAY, 0.5)
        tb(slide, x+0.05, y+0.02, w-0.1, 0.36, val, 11, False, TEXT_PRIMARY, PP_ALIGN.CENTER)

# 右侧：各产线核心指标确认
tb(slide, 9.5, 1.2, 3.5, 0.4, "各产线核心指标", 17, True, ACCENT_TEAL)
tb(slide, 9.5, 1.6, 3.5, 0.3, "每产线确认2-3个核心前置指标", 12, False, TEXT_SECONDARY)

team_list = [("无线", ACCENT_BLUE), ("交换机", ACCENT_PURPLE), ("NMC", ACCENT_TEAL),
             ("会议主机", ACCENT_ORANGE), ("IPSIP", RGBColor(0x7B, 0x83, 0xEB)),
             ("硬件", RGBColor(0xE9, 0x7C, 0x23))]
for i, (team, c) in enumerate(team_list):
    y = 2.1 + i * 0.95
    rrect(slide, 9.5, y, 3.3, 0.8, TEXT_WHITE, c, 1)
    tb(slide, 9.7, y+0.05, 1.2, 0.3, team, 13, True, c)
    tb(slide, 10.9, y+0.05, 1.7, 0.7, "指标①：\n指标②：", 11, False, TEXT_SECONDARY)

# 底部提示
rrect(slide, 0.8, 6.9, 11.5, 0.4, BG_LIGHT_BLUE, ACCENT_TEAL, 1)
tb(slide, 1, 6.93, 11, 0.34, "💡 核心原则：前置指标不必多，2-3个能跟踪就好，后续迭代优化", 13, True, ACCENT_TEAL, PP_ALIGN.CENTER)

# --- 检视节奏 ---
slide = content_slide("环节三：检视节奏")

tb(slide, 0.8, 1.3, 11, 0.4, "建议的检视节奏，请确认或调整：", 16, False, TEXT_SECONDARY)

cadences = [
    ("每周", "任务执行状态更新\n障碍状态更新", "产线内部", "进度同步", ACCENT_TEAL),
    ("双周", "前置指标追踪\n障碍闭环情况", "产线 + 综合管理部", "红黄绿灯报告", ACCENT_BLUE),
    ("每月", "里程碑达成评估\n目标偏差分析", "全体", "月度检视报告", ACCENT_ORANGE),
    ("每季度", "OKR达成评估\n里程碑复盘\n下季度目标调整", "全体 + 管理层", "季度复盘报告", ACCENT_PURPLE),
]

for i, (freq, content, participants, output, color) in enumerate(cadences):
    y = 1.9 + i * 1.3
    # 频率标签
    rrect(slide, 0.8, y, 1.5, 1.1, color)
    tb(slide, 0.8, y+0.3, 1.5, 0.5, freq, 20, True, TEXT_WHITE, PP_ALIGN.CENTER)
    # 内容
    rrect(slide, 2.5, y, 3.5, 1.1, TEXT_WHITE, BORDER_GRAY, 1)
    tb(slide, 2.7, y+0.1, 3.2, 0.9, content, 13, False, TEXT_PRIMARY)
    # 参与人
    tb(slide, 6.3, y+0.1, 2.5, 0.9, participants, 13, False, TEXT_SECONDARY)
    # 产出
    rrect(slide, 9, y+0.1, 3.5, 0.9, BG_LIGHT_BLUE, color, 1)
    tb(slide, 9.2, y+0.2, 3.2, 0.7, output, 14, True, color, PP_ALIGN.CENTER)

# --- 红黄绿灯 ---
slide = content_slide("环节三：红黄绿灯报告机制")

tb(slide, 0.8, 1.3, 11, 0.4, "双周检视使用红黄绿灯报告，快速暴露偏差：", 16, False, TEXT_SECONDARY)

lights = [
    ("🟢 正常", "按计划推进", "不需要额外关注\n下次检视继续跟踪", ACCENT_GREEN),
    ("🟡 有风险", "有偏差但可控", "双周检视中讨论\n寻找解决方案", ACCENT_ORANGE),
    ("🔴 卡住", "无法推进，需帮助", "立即升级\n下次检视前必须有行动", ACCENT_RED),
]

for i, (status, desc, action, color) in enumerate(lights):
    x = 1 + i * 4
    rrect(slide, x, 2, 3.5, 3.5, TEXT_WHITE, color, 2)
    tb(slide, x+0.3, 2.2, 2.9, 0.6, status, 28, True, color, PP_ALIGN.CENTER)
    rrect(slide, x+0.3, 3, 2.9, 0.6, color)
    tb(slide, x+0.3, 3.1, 2.9, 0.4, desc, 16, True, TEXT_WHITE, PP_ALIGN.CENTER)
    tb(slide, x+0.3, 3.9, 2.9, 1.2, action, 14, False, TEXT_SECONDARY, PP_ALIGN.CENTER)

# 报告模板
tb(slide, 0.8, 5.9, 11, 0.4, "报告模板：", 15, True, TEXT_PRIMARY)

rpt_cols = [("里程碑/目标", 3), ("状态", 1.2), ("偏差说明", 3), ("需要的支持", 3.2)]
x = 0.8
for name, w in rpt_cols:
    rect(slide, x, 6.35, w, 0.45, RGBColor(0x33, 0x33, 0x33))
    tb(slide, x, 6.38, w, 0.4, name, 13, True, TEXT_WHITE, PP_ALIGN.CENTER)
    x += w

# 示例行
rpt_row = ["投屏模块CI门禁", "🟡", "代码审查skill未完成", "AI效能组支持skill开发"]
x = 0.8
for i, (val, (_, w)) in enumerate(zip(rpt_row, rpt_cols)):
    rect(slide, x, 6.8, w, 0.4, BG_ROW_ALT, BORDER_GRAY, 0.5)
    tb(slide, x+0.05, 6.82, w-0.1, 0.36, val, 12, False, TEXT_PRIMARY, PP_ALIGN.CENTER)
    x += w

# --- 行动项确认 ---
slide = content_slide("环节三：行动项确认与收尾")

tb(slide, 0.8, 1.3, 11, 0.4, "确认本次研讨会的产出和行动项：", 16, False, TEXT_SECONDARY)

# 产出回顾
rrect(slide, 0.8, 1.9, 3.5, 1.5, TEXT_WHITE, ACCENT_RED, 1.5)
tb(slide, 1, 2, 3, 0.35, "产出① 障碍清单", 15, True, ACCENT_RED)
tb(slide, 1, 2.4, 3, 0.8, "分类 · 优先级\n解法方向 · 责任方", 13, False, TEXT_PRIMARY)

rrect(slide, 4.8, 1.9, 3.5, 1.5, TEXT_WHITE, ACCENT_BLUE, 1.5)
tb(slide, 5, 2, 3, 0.35, "产出② 目标拆解", 15, True, ACCENT_BLUE)
tb(slide, 5, 2.4, 3, 0.8, "SMART验证 · 行动项\n责任人 · 时间", 13, False, TEXT_PRIMARY)

rrect(slide, 8.8, 1.9, 3.5, 1.5, TEXT_WHITE, ACCENT_GREEN, 1.5)
tb(slide, 9, 2, 3, 0.35, "产出③ 检视机制", 15, True, ACCENT_GREEN)
tb(slide, 9, 2.4, 3, 0.8, "前置指标 · 红黄绿灯\n检视节奏", 13, False, TEXT_PRIMARY)

# 会后跟进
tb(slide, 0.8, 3.8, 11, 0.4, "会后跟进时间线：", 17, True, TEXT_PRIMARY)

followups = [
    ("会后1天", "发布会议纪要 + 行动项清单", "主持人", ACCENT_BLUE),
    ("会后1周", "各产线补充完成全部SMART拆解", "各产线", ACCENT_PURPLE),
    ("会后2周", "第一次双周检视（红黄绿灯报告）", "各产线", ACCENT_ORANGE),
    ("会后1月", "第一次月度检视", "全体", ACCENT_TEAL),
]

for i, (time, action, who, c) in enumerate(followups):
    y = 4.4 + i * 0.65
    rrect(slide, 0.8, y, 1.8, 0.5, c)
    tb(slide, 0.8, y+0.05, 1.8, 0.4, time, 14, True, TEXT_WHITE, PP_ALIGN.CENTER)
    rrect(slide, 2.8, y, 6, 0.5, TEXT_WHITE, BORDER_GRAY, 1)
    tb(slide, 3, y+0.05, 5.6, 0.4, action, 14, False, TEXT_PRIMARY)
    tb(slide, 9, y+0.05, 3, 0.4, who, 13, False, TEXT_SECONDARY, PP_ALIGN.CENTER)

# --- 结尾页 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
rect(slide, 1, 3.3, 2, 0.06, ACCENT_BLUE)
tb(slide, 1, 2.2, 11, 1, "研讨会到此结束", 36, True, TEXT_WHITE)
tb(slide, 1, 3.6, 11, 0.6, "会后1周内完成补充拆解，2周后第一次双周检视", 20, False, RGBColor(0xBB, 0xBB, 0xBB))
tb(slide, 1, 5.5, 11, 0.4, "感谢各位参与，行动起来 💪", 18, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════
output = "/root/.openclaw/workspace/knowledge/work/AI-Native/AI-Native落地推进研讨会-开场共识引导.pptx"
prs.save(output)
print(f"PPT已生成：{output}")
print(f"共 {len(prs.slides)} 页")
